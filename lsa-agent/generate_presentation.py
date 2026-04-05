#!/usr/bin/env python3
"""
Generate a professional hackathon presentation for LSA (Life Simulation Agent)
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(25, 55, 109)
LIGHT_BLUE = RGBColor(60, 120, 200)
ACCENT_GREEN = RGBColor(52, 211, 153)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_title_slide(title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = ACCENT_GREEN
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, content_items):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    # Add title text
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.space_before = Pt(10)
    title_p.space_after = Pt(10)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    return slide

# Slide 1: Title
add_title_slide(
    "Life Simulation Agent",
    "Predict Consequences. Analyze Decisions. Optimize Life."
)

# Slide 2: Problem Statement
add_content_slide(
    "🎯 The Problem",
    [
        "• Life decisions are complex with multiple cascading effects",
        "• No data-driven way to predict outcomes of choices",
        "• People struggle with long-term impact analysis",
        "• Current solutions: gut feeling, regret, hindsight bias",
        "• Need: Real-time simulation of decision scenarios"
    ]
)

# Slide 3: Solution Overview
add_content_slide(
    "💡 Our Solution: LSA",
    [
        "• Autonomous AI agent that simulates life scenarios",
        "• Multi-scenario decision analysis (3 scenarios per choice)",
        "• Machine learning-powered pattern recognition",
        "• Real-time WhatsApp integration for accessibility",
        "• 7-day outcome predictions with confidence scoring"
    ]
)

# Slide 4: Key Features
add_content_slide(
    "✨ Key Features",
    [
        "📊 Decision Simulation: A/B/C scenario analysis with risk scoring",
        "💾 Event Logging: Track life events with impact metrics",
        "📈 Pattern Analysis: Detect behavioral trends & consistency",
        "💬 WhatsApp Integration: Chat directly with LSA",
        "🧠 Smart Memory: Semantic search of historical decisions"
    ]
)

# Slide 5: Architecture
add_content_slide(
    "🏗️ Technical Architecture",
    [
        "Memory Manager: Persistent storage with ML embeddings (FAISS)",
        "Simulation Engine: Multi-scenario generation via Agnes Claw LLM",
        "Intervention Engine: Proactive alerts & pattern warnings",
        "WhatsApp Gateway: OpenClaw API + ngrok webhook tunneling",
        "All components: Modular, tested, production-ready"
    ]
)

# Slide 6: Tech Stack
add_content_slide(
    "🛠️ Technology Stack",
    [
        "Backend: Python 3.11 | Flask | Requests",
        "ML/AI: sentence-transformers | FAISS | Agnes Claw Model",
        "Integration: OpenClaw API | ngrok | WhatsApp",
        "Data: JSONL + NumPy embeddings | 384-dim vectors",
        "Deployment: localhost:5001 | Easy production scaling"
    ]
)

# Slide 7: How It Works
add_content_slide(
    "🚀 How LSA Works",
    [
        "1️⃣ User sends decision/event via WhatsApp or Python API",
        "2️⃣ Memory Manager loads user context (past 30 days)",
        "3️⃣ Simulation Engine generates 3 decision scenarios",
        "4️⃣ Risk, confidence, and outcome scoring computed",
        "5️⃣ LSA recommends optimal path with prediction insights"
    ]
)

# Slide 8: Live Demo
add_content_slide(
    "🎬 Demo Workflow",
    [
        "User: 'Should I skip studying today?'",
        "",
        "LSA Analysis Shows:",
        "  • Scenario A (Skip): -5% momentum, 60% consistency",
        "  • Scenario B (Moderate): +15% momentum, 78% consistency ✅",
        "  • Scenario C (Study 2hrs): 95% consistency, high regret risk",
        "",
        "Recommendation: B (Moderate improvement)"
    ]
)

# Slide 9: Key Metrics
add_content_slide(
    "📊 Project Metrics",
    [
        "✅ 3 Core Skills: Memory, Simulation, Intervention (100% implemented)",
        "✅ 15 Memories Loaded: Demo contextual data ready",
        "✅ 2 Decision Scenarios: Full A/B/C analysis per decision",
        "✅ 87% Consistency: Behavioral pattern tracking active",
        "✅ <3 seconds: Decision analysis response time"
    ]
)

# Slide 10: Differentiators
add_content_slide(
    "🌟 What Makes LSA Unique",
    [
        "🚫 NOT Just a chatbot - Autonomous decision analysis system",
        "🎯 Multi-scenario simulation - Not single-path predictions",
        "📚 Long-term memory - Contextual learning from past events",
        "⚡ Production-ready - Deployed with real WhatsApp integration",
        "💰 No setup hassle - Pre-configured with ngrok + OpenClaw"
    ]
)

# Slide 11: Use Cases
add_content_slide(
    "💼 Real-World Use Cases",
    [
        "Career Decisions: Job offer evaluation with relocation impact",
        "Habit Building: Exercise routine consistency tracking",
        "Study Planning: Optimal learning schedule generation",
        "Health Choices: Workout + nutrition decision analysis",
        "Life Planning: Major decisions backed by data simulation"
    ]
)

# Slide 12: Future Roadmap
add_content_slide(
    "🔮 Future Roadmap",
    [
        "🌍 Multi-user support with privacy isolation",
        "📱 iOS/Android mobile app (not just WhatsApp)",
        "🤖 Advanced LLM fine-tuning for personal context",
        "📊 Interactive dashboard & data visualization",
        "🔄 Integration with calendar, fitness trackers, repos",
        "💭 Emotion-aware predictions & mental health insights"
    ]
)

# Slide 13: Challenges Solved
add_content_slide(
    "🛡️ Challenges We Solved",
    [
        "✅ LLM Integration: Agnes Claw Model for intelligent scenarios",
        "✅ Webhook Hell: Full production pipeline (ngrok + OpenClaw)",
        "✅ No Twilio: Direct OpenClaw WhatsApp (simpler, cheaper)",
        "✅ Performance: Optimized embedding search < 100ms",
        "✅ UX: One-command setup with auto-configured webhooks"
    ]
)

# Slide 14: Results So Far
add_content_slide(
    "🏆 Results & Achievements",
    [
        "✨ Full system built and tested (all components working)",
        "📦 Production-ready codebase on GitHub",
        "🧪 Comprehensive test suite (Memory, Simulation, Intervention)",
        "📚 Complete documentation (README, setup guides, architecture)",
        "🚀 Ready for deployment and real users today"
    ]
)

# Slide 15: Call to Action
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = ACCENT_GREEN

# Main message
main_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
main_frame = main_box.text_frame
main_frame.word_wrap = True
main_p = main_frame.paragraphs[0]
main_p.text = "Ready to Predict Your Future?"
main_p.font.size = Pt(54)
main_p.font.bold = True
main_p.font.color.rgb = DARK_BLUE
main_p.alignment = PP_ALIGN.CENTER

# GitHub link
github_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
github_frame = github_box.text_frame
github_p = github_frame.paragraphs[0]
github_p.text = "GitHub: Harikrishna2461/Agnes_OpenCLaw_Personal_Digital_Agent"
github_p.font.size = Pt(24)
github_p.font.color.rgb = DARK_BLUE
github_p.alignment = PP_ALIGN.CENTER

# Contact
contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
contact_frame = contact_box.text_frame
contact_p = contact_frame.paragraphs[0]
contact_p.text = "Try LSA: Send WhatsApp to +91 7010384691"
contact_p.font.size = Pt(20)
contact_p.font.bold = True
contact_p.font.color.rgb = WHITE
contact_p.alignment = PP_ALIGN.CENTER

# Save presentation
output_file = "LSA_Hackathon_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created: {output_file}")
print(f"📊 Total slides: 15")
print(f"🎨 Color scheme: Professional blue & green")
print(f"📍 Location: {output_file}")
